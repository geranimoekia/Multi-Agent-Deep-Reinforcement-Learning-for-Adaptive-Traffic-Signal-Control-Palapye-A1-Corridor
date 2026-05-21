from pptx import Presentation
from pptx.util import Inches

prs = Presentation('docs/TSOTLHE SEIPHEPI FYP.pptx')
print(f"Size: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")
print(f"Layouts: {len(prs.slide_layouts)}")

for i, sl in enumerate(prs.slides):
    print(f'\n--- Slide {i+1} ---')
    for sh in sl.shapes:
        try:
            fill = sh.fill
            ft = fill.type
            try:
                color = str(fill.fore_color.rgb)
            except Exception:
                color = f'type={ft}'
        except Exception:
            color = 'err'
        try:
            txt = sh.text[:60].replace('\n', ' ') if hasattr(sh, 'text') else ''
        except Exception:
            txt = ''
        l = round(sh.left/914400, 2)
        t = round(sh.top/914400, 2)
        w = round(sh.width/914400, 2)
        h = round(sh.height/914400, 2)
        print(f'  {sh.name:28s} pos=({l},{t}) sz=({w}x{h}) fill={color} txt={txt!r}')

# Also check master background
sm = prs.slide_masters[0]
print("\n--- Slide Master shapes ---")
for sh in sm.shapes:
    try:
        txt = sh.text[:60].replace('\n', ' ') if hasattr(sh, 'text') else ''
    except Exception:
        txt = ''
    l = round(sh.left/914400, 2)
    t = round(sh.top/914400, 2)
    w = round(sh.width/914400, 2)
    h = round(sh.height/914400, 2)
    print(f'  {sh.name:28s} pos=({l},{t}) sz=({w}x{h}) txt={txt!r}')
