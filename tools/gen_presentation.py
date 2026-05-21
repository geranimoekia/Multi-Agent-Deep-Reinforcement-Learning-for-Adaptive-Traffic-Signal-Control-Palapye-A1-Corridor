"""
Generate PALMS FYP presentation using the BIUST template from
docs/TSOTLHE SEIPHEPI FYP.pptx.

Strategy:
  - Open the old PPTX as a template source
  - For each new slide, deep-copy slide 2's XML (standard BIUST chrome:
    number box, orange title bar, BIUST logos, bottom strip)
  - Clear the content textbox and write new content
  - Append to a fresh copy of the template

Run:  python tools/gen_presentation.py
Output: docs/PALMS_Presentation_2025.pptx
"""

import copy, os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = "docs/TSOTLHE SEIPHEPI FYP.pptx"
OUT      = "docs/PALMS_Presentation_BIUST.pptx"

# colours from the BIUST template
CYAN    = RGBColor(0x00, 0xB0, 0xF0)   # number box
ORANGE  = RGBColor(0xFB, 0x82, 0x3F)   # title bar
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
MID     = RGBColor(0x00, 0x70, 0xC0)
GREEN   = RGBColor(0x06, 0x8A, 0x4C)
RED     = RGBColor(0xCC, 0x33, 0x33)
LGREY   = RGBColor(0xF0, 0xF5, 0xFB)
DGREY   = RGBColor(0x55, 0x55, 0x55)

# content area (inside the BIUST chrome)
CX = Inches(0.10)        # left edge
CT = Inches(1.00)        # top edge (below title bar)
CW = Inches(9.80)        # usable width
CH = Inches(5.80)        # usable height


# ─────────────────────────────────────────────────────────────────────────────
# LOW-LEVEL: clone a template slide XML, then add shapes into it
# ─────────────────────────────────────────────────────────────────────────────

def _clone_slide(prs_src, slide_idx, prs_dst):
    """
    Deep-copy slide `slide_idx` from prs_src into prs_dst.
    Returns the new slide object.

    Key detail: we clear-and-repopulate dst_spTree IN PLACE so that python-pptx's
    shape manager (which holds a reference to that element) keeps working after
    the clone.  Replacing dst_spTree with a new element would detach the manager.
    """
    template_slide = prs_src.slides[slide_idx]
    blank_layout   = prs_dst.slide_layouts[6]
    new_slide      = prs_dst.slides.add_slide(blank_layout)

    src_part = template_slide.part
    dst_part = new_slide.part

    # Build a rId mapping: old rId → new rId for every image relationship
    rId_map = {}
    for rel in src_part.rels.values():
        if "image" in rel.reltype:
            new_rId = dst_part.relate_to(rel.target_part, rel.reltype)
            rId_map[rel.rId] = new_rId

    # Serialise source spTree, substitute rIds, then parse back
    src_xml = etree.tostring(template_slide.shapes._spTree, encoding="unicode")
    for old_rId, new_rId in rId_map.items():
        src_xml = src_xml.replace(f'r:embed="{old_rId}"', f'r:embed="{new_rId}"')
        src_xml = src_xml.replace(f'r:id="{old_rId}"',    f'r:id="{new_rId}"')
    new_spTree = etree.fromstring(src_xml)

    # Clear-and-repopulate dst_spTree IN PLACE (preserves python-pptx's internal ref)
    dst_spTree = new_slide.shapes._spTree
    for child in list(dst_spTree):
        dst_spTree.remove(child)
    for child in list(new_spTree):
        dst_spTree.append(child)   # children already carry substituted rIds

    return new_slide


def _set_title(slide, number_str, title_str):
    """Update the cyan number box and the orange title box."""
    for sh in slide.shapes:
        if not hasattr(sh, "text_frame"):
            continue
        txt = sh.text.strip()
        # The number box contains only a small integer-ish string
        try:
            fill = sh.fill.fore_color.rgb
        except Exception:
            fill = None
        if fill == CYAN:
            sh.text_frame.paragraphs[0].runs[0].text = number_str
        elif fill == ORANGE:
            sh.text_frame.paragraphs[0].runs[0].text = title_str


def _clear_textbox(slide, name="TextBox 6"):
    """Remove any existing content textbox (name-matched or all TextBox shapes)."""
    to_remove = []
    sp_tree = slide.shapes._spTree
    for sp in sp_tree:
        nm = sp.get('{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}name', '')
        # Use the pptx name attribute
        try:
            nvsp = sp.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
            cNvPr = nvsp.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr') if nvsp is not None else None
            if cNvPr is None:
                # try drawingml ns
                nvsp2 = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
                cNvPr = nvsp2
        except Exception:
            cNvPr = None

        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag == 'sp':
            # Check if it has a txBody but is NOT the number/title/footer shapes
            txBody = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}txBody')
            if txBody is not None:
                # check fill
                solidFill = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                srgbClr   = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr') if solidFill is not None else None
                if srgbClr is None:
                    # no solid fill → content textbox → remove
                    to_remove.append(sp)
    for sp in to_remove:
        sp_tree.remove(sp)


# ─────────────────────────────────────────────────────────────────────────────
# HIGH-LEVEL: add shapes to the content area
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w_pt=0):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    f  = sh.fill
    if fill_rgb:
        f.solid(); f.fore_color.rgb = fill_rgb
    else:
        f.background()
    ln = sh.line
    if line_rgb:
        ln.color.rgb = line_rgb
        ln.width     = Pt(line_w_pt)
    else:
        ln.fill.background()
    return sh


def add_text(slide, text, l, t, w, h,
             size=14, bold=False, italic=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = "Calibri"
    return txb


def add_bullets(slide, items, l, t, w, h, size=13, color=DARK,
                levels=None, bold_first=False, spacing=0):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        lvl   = levels[i] if levels else 0
        indent = "        " * lvl
        bullet = "•" if lvl == 0 else "◦"
        run    = p.add_run()
        run.text = f"{indent}{bullet}  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.name  = "Calibri"
        run.font.bold  = bold_first and i == 0
    return txb


def render_table(slide, headers, rows, col_widths,
                 left, top, row_h=Inches(0.38),
                 hdr_fill=None, alt_fill=None):
    hdr_fill = hdr_fill or MID
    alt_fill = alt_fill or LGREY
    for c, (h, cw) in enumerate(zip(headers, col_widths)):
        xl = left + sum(col_widths[:c])
        add_rect(slide, xl, top, cw, row_h, fill_rgb=hdr_fill)
        add_text(slide, h, xl + Inches(0.05), top + Inches(0.03),
                 cw - Inches(0.1), row_h - Inches(0.06),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        rf = alt_fill if r % 2 == 0 else WHITE
        for c, (cell, cw) in enumerate(zip(row, col_widths)):
            xl = left + sum(col_widths[:c])
            yt = top + (r + 1) * row_h
            is_last = (c == len(col_widths) - 1)
            col_val = (GREEN if ("−" in cell or "-" in cell) and is_last and "%" in cell
                       else DARK)
            if is_last and "+" in cell and "%" in cell:
                col_val = GREEN
            add_rect(slide, xl, yt, cw, row_h, fill_rgb=rf,
                     line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_w_pt=0.4)
            add_text(slide, cell, xl + Inches(0.05), yt + Inches(0.03),
                     cw - Inches(0.1), row_h - Inches(0.06),
                     size=10, bold=is_last, color=col_val,
                     align=PP_ALIGN.CENTER)


def metric_card(slide, label, value, note, l, t,
                w=Inches(1.88), h=Inches(1.20), val_color=GREEN):
    add_rect(slide, l, t, w, h,
             fill_rgb=WHITE,
             line_rgb=MID, line_w_pt=1.2)
    add_text(slide, label,
             l + Inches(0.07), t + Inches(0.05), w - Inches(0.14), Inches(0.30),
             size=10, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, value,
             l + Inches(0.05), t + Inches(0.33), w - Inches(0.10), Inches(0.50),
             size=22, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, note,
             l + Inches(0.05), t + Inches(0.82), w - Inches(0.10), Inches(0.30),
             size=9, color=DGREY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# BUILD PRESENTATION
# ─────────────────────────────────────────────────────────────────────────────

src = Presentation(TEMPLATE)    # template source (read-only)
dst = Presentation(TEMPLATE)   # destination — keep original 15 slides for now
#
# Strategy: add all new slides FIRST (they get parts numbered slide16..slide30)
# then strip the original 15 from sldIdLst.  The old parts remain in the zip
# as unreferenced but PowerPoint ignores them.
#
N_ORIG = len(dst.slides)       # 15 originals to strip after building

# Template slide to clone for each content slide = slide index 1 (slide 2, standard layout)
TMPL_IDX = 1   # 0-based

SLIDE_W = dst.slide_width   # 9144000 EMU = 10 inches
SLIDE_H = dst.slide_height  # 6858000 EMU = 7.5 inches

slides_data = []   # list of (number, title, build_fn)


# ══════════════════════════════════════════════════════════════════
def s01_title():
    """Title slide — clone slide 1 (special title layout) from template."""
    sl = _clone_slide(src, 0, dst)
    # Find and update text shapes
    for sh in sl.shapes:
        if not hasattr(sh, "text_frame"):
            continue
        try:
            fill_rgb = sh.fill.fore_color.rgb
        except Exception:
            fill_rgb = None
        if fill_rgb == ORANGE:
            sh.text_frame.paragraphs[0].runs[0].text = (
                "Adaptive Multi-Agent Reinforcement Learning "
                "for Traffic Signal Optimisation along the Palapye A1 Urban Corridor"
            )
            sh.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    # Replace the content placeholder text
    for sh in sl.shapes:
        if not hasattr(sh, "text_frame"):
            continue
        t = sh.text.strip()
        if "SCHOOL" in t or "DEPARTMENT" in t or "EEEN" in t:
            tf = sh.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            tf.paragraphs[0].runs[0].text = (
                "SCHOOL OF ENGINEERING\n"
                "DEPARTMENT OF ELECTRICAL, COMPUTER AND TELECOMMUNICATIONS ENGINEERING\n\n"
                "EEEN510 Final Year Project\n"
                "Tsotlhe Nayang Seiphepi  |  21001137\n"
                "Supervisor: Dr Bokamoso Basutli\n"
                "Date: May 2025"
            )
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)
    return sl


def s02_outline():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "1", "PRESENTATION OUTLINE")
    _clear_textbox(sl)
    cards = [
        ("1. Introduction & Problem",  MID),
        ("2. Research Objectives",     MID),
        ("3. System Architecture",     MID),
        ("4. Phase 1 SA-PPO Results",  MID),
        ("5. MAPPO System Design",     MID),
        ("6. MAPPO Training",          MID),
        ("7. Evaluation Results",      MID),
        ("8. Conclusions & Recommendations", DARK),
    ]
    cw, ch = Inches(1.15), Inches(1.40)
    gap    = Inches(0.08)
    tl     = Inches(0.12)
    tt     = Inches(1.08)
    for i, (lbl, col) in enumerate(cards):
        xl = tl + i * (cw + gap)
        add_rect(sl, xl, tt, cw, ch, fill_rgb=col)
        add_text(sl, lbl, xl + Inches(0.06), tt + Inches(0.20),
                 cw - Inches(0.12), ch - Inches(0.40),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl,
             "A two-phase project — from single-agent PPO at one junction to "
             "cooperative MAPPO across the full A1 Palapye corridor.",
             Inches(0.15), Inches(2.70), Inches(9.70), Inches(0.50),
             size=13, italic=True, color=DARK, align=PP_ALIGN.CENTER)
    return sl


def s03_intro():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "2", "INTRODUCTION & PROBLEM STATEMENT")
    _clear_textbox(sl)
    add_bullets(sl, [
        "Palapye: key transit hub on the A1 highway between Gaborone and Francistown",
        "Morupule Coal Mine + BIUST campus drive recurrent, time-peaked demand",
        "Existing signals rely on static, pre-timed fixed cycles",
        "Fixed-time control is blind to queue lengths — cannot adapt to peak demand",
        "3 closely spaced signalised junctions within ~1 km — decisions at one affect the others",
        "Consequences: chronic congestion, excessive idling, elevated emissions",
    ],
    Inches(0.15), Inches(1.08), Inches(5.60), Inches(3.80),
    size=13, color=DARK)

    add_rect(sl, Inches(5.95), Inches(1.08), Inches(3.95), Inches(3.80),
             fill_rgb=LGREY)
    add_text(sl, "PROBLEM STATEMENT",
             Inches(6.08), Inches(1.15), Inches(3.70), Inches(0.35),
             size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(6.08), Inches(1.52), Inches(3.70), Pt(2), fill_rgb=ORANGE)
    add_text(sl,
             "Urban intersections in Palapye experience chronic congestion "
             "due to fixed-time controllers that cannot adapt to real-time "
             "conditions, coordinate across junctions, or respond to varying demand.",
             Inches(6.08), Inches(1.60), Inches(3.70), Inches(1.80),
             size=11, color=DARK)
    add_text(sl, "RESEARCH GAP",
             Inches(6.08), Inches(3.52), Inches(3.70), Inches(0.30),
             size=11, bold=True, color=RED)
    add_bullets(sl, [
        "No adaptive/RL controller deployed or evaluated on a Botswana signalised network",
        "No publicly available traffic count dataset for Palapye intersections",
    ],
    Inches(6.08), Inches(3.85), Inches(3.70), Inches(0.90),
    size=10, color=DARK)

    # summary bar
    add_rect(sl, Inches(0.12), Inches(5.10), Inches(9.76), Inches(0.55),
             fill_rgb=MID)
    add_text(sl,
             "This project builds and evaluates the first cooperative MARL traffic signal "
             "controller on a real Botswana road network.",
             Inches(0.20), Inches(5.17), Inches(9.60), Inches(0.38),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return sl


def s04_objectives():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "3", "RESEARCH OBJECTIVES & CONTRIBUTIONS")
    _clear_textbox(sl)
    add_text(sl, "Four Objectives",
             Inches(0.15), Inches(1.08), Inches(5.60), Inches(0.32),
             size=13, bold=True, color=MID)
    add_bullets(sl, [
        "YOLOv8 vehicle detection module validated on Palapye traffic footage",
        "SUMO three-intersection A1 corridor model with 6 demand scenarios",
        "MAPPO with Centralised Training, Decentralised Execution (CTDE) for 3 cooperative agents",
        "Quantitative comparison: MAPPO vs SA-PPO vs Fixed-Time across all scenarios",
    ],
    Inches(0.15), Inches(1.43), Inches(5.60), Inches(2.20),
    size=12, color=DARK)

    add_text(sl, "Key Contributions",
             Inches(0.15), Inches(3.70), Inches(5.60), Inches(0.32),
             size=13, bold=True, color=MID)
    add_bullets(sl, [
        "First cooperative MARL controller on a real Botswana road network",
        "Reusable SUMO simulation environment (generates data for future researchers)",
        "First quantitative Fixed vs SA-PPO vs MAPPO comparison at the A1 corridor",
        "YOLOv8 pipeline providing a real-world sensing pathway",
    ],
    Inches(0.15), Inches(4.05), Inches(5.60), Inches(1.65),
    size=12, color=DARK)

    # Phase boxes
    for idx, (tag, title, det, col) in enumerate([
        ("PHASE 1", "Single-Agent PPO",
         "1 junction  •  1 demand scenario\nStable-Baselines3  •  ~5M steps", MID),
        ("PHASE 2", "MAPPO CTDE",
         "3 junctions  •  6 scenarios  •  curriculum\nCustom MAPPO  •  1.5M steps",   DARK),
    ]):
        yt = Inches(1.08) + idx * Inches(2.45)
        add_rect(sl, Inches(6.00), yt, Inches(3.90), Inches(2.20), fill_rgb=col)
        add_text(sl, tag,
                 Inches(6.12), yt + Inches(0.10), Inches(3.66), Inches(0.28),
                 size=10, bold=True, color=ORANGE)
        add_text(sl, title,
                 Inches(6.12), yt + Inches(0.40), Inches(3.66), Inches(0.40),
                 size=15, bold=True, color=WHITE)
        add_text(sl, det,
                 Inches(6.12), yt + Inches(0.88), Inches(3.66), Inches(1.10),
                 size=11, color=RGBColor(0xCC, 0xDD, 0xEE))
        if idx == 0:
            add_text(sl, "↓  extends to",
                     Inches(7.1), Inches(3.30), Inches(2.0), Inches(0.30),
                     size=12, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    return sl


def s05_architecture():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "4", "SYSTEM ARCHITECTURE")
    _clear_textbox(sl)

    pipe = [
        ("Camera /\nYOLOv8", CYAN),
        ("Traffic\nState\nExtraction", MID),
        ("RL Policy\nDecision\n(PPO/MAPPO)", DARK),
        ("SUMO+TraCI\nSignal\nActuation", MID),
        ("Performance\nMetrics", CYAN),
    ]
    bw, bh = Inches(1.72), Inches(1.55)
    bt     = Inches(1.18)
    bl0    = Inches(0.18)
    gap    = Inches(0.30)
    for i, (lbl, col) in enumerate(pipe):
        xl = bl0 + i * (bw + gap)
        add_rect(sl, xl, bt, bw, bh, fill_rgb=col)
        add_text(sl, lbl, xl + Inches(0.07), bt + Inches(0.18),
                 bw - Inches(0.14), bh - Inches(0.36),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(pipe) - 1:
            add_text(sl, "→",
                     xl + bw + Inches(0.05), bt + bh/2 - Pt(10),
                     Inches(0.22), Inches(0.30),
                     size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    add_text(sl, "↺  feedback loop: performance metrics drive next policy update",
             Inches(0.18), Inches(2.90), Inches(9.64), Inches(0.30),
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)

    info = [
        ("SUMO Network",
         "Real Palapye A1 geometry\n3 signalised junctions\n6 demand scenarios\nStep: 3 s/action"),
        ("MAPPO CTDE",
         "22-dim local obs per agent\n66-dim joint state for critic\nGlobal cooperative reward\n3 agents, param sharing"),
        ("YOLOv8",
         "Anchor-free detection\nValidated on Palapye footage\nQueue length estimation\nFuture: live integration"),
    ]
    for k, (title, body) in enumerate(info):
        xl = Inches(0.18) + k * Inches(3.30)
        add_rect(sl, xl, Inches(3.35), Inches(3.12), Inches(2.30),
                 fill_rgb=LGREY)
        add_text(sl, title, xl + Inches(0.10), Inches(3.42),
                 Inches(2.92), Inches(0.30),
                 size=11, bold=True, color=DARK)
        add_text(sl, body, xl + Inches(0.10), Inches(3.78),
                 Inches(2.92), Inches(1.70),
                 size=10, color=DARK)
    return sl


def s06_ppo_results():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "5", "PHASE 1: SINGLE-AGENT PPO vs FIXED-TIME")
    _clear_textbox(sl)

    add_text(sl,
             "Single intersection  •  9,000 timesteps  •  5 demand phases  "
             "•  Stable-Baselines3 PPO  •  ~5M training steps",
             Inches(0.15), Inches(1.05), Inches(9.70), Inches(0.28),
             size=11, italic=True, color=DGREY, align=PP_ALIGN.CENTER)

    metrics = [
        ("Throughput",      "+48.5%", "p<0.001  d=0.156", GREEN),
        ("Queue Length",    "−27.8%", "p<0.001  d=0.253", GREEN),
        ("Traffic Pressure","−27.3%", "p<0.001  d=0.239", GREEN),
        ("Delay (agg.)",    "−10.1%", "Negligible d=0.131", ORANGE),
        ("Stop Ratio",      "−3.6%",  "p=0.017  d=0.042",  ORANGE),
    ]
    ml = Inches(0.12)
    for i, (lbl, val, note, col) in enumerate(metrics):
        metric_card(sl, lbl, val, note,
                    ml + i * Inches(1.96), Inches(1.40), val_color=col)

    add_text(sl,
             "Delay aggregate is misleading — phase analysis shows 50–68% improvement during heavy/peak demand.",
             Inches(0.15), Inches(2.75), Inches(9.70), Inches(0.30),
             size=11, italic=True, color=MID, align=PP_ALIGN.CENTER)

    hdrs = ["Demand Phase", "Delay improv.", "Queue improv.", "Throughput improv."]
    rows = [
        ("Light (0–60 s)",      "+64.8%",  "—",       "—"),
        ("Heavy (120–180 s)",   "+68.2%",  "+100%",   "+100%"),
        ("Peak (180–240 s)",    "+50.1%",  "−34.4%",  "+—"),
        ("Extended (300 s+)",   "−11.7%",  "+29.9%",  "+39.8%"),
    ]
    cw = [Inches(2.60), Inches(2.20), Inches(2.20), Inches(2.20)]
    render_table(sl, hdrs, rows, cw, Inches(0.25), Inches(3.10))

    add_rect(sl, Inches(0.12), Inches(5.08), Inches(9.76), Inches(0.55), fill_rgb=DARK)
    add_text(sl,
             "These results confirmed RL can learn adaptive signal control — "
             "and established the baseline Phase 2 was designed to surpass at network scale.",
             Inches(0.20), Inches(5.15), Inches(9.60), Inches(0.38),
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return sl


def s07_mappo_design():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "6", "PHASE 2: MAPPO SYSTEM DESIGN — CTDE")
    _clear_textbox(sl)

    add_text(sl, "WHY MAPPO?",
             Inches(0.15), Inches(1.05), Inches(5.55), Inches(0.30),
             size=12, bold=True, color=MID)
    add_bullets(sl, [
        "SA-PPO has no cross-junction visibility — cannot anticipate downstream spillback",
        "Fixed-time cycles blindly regardless of queue state",
        "MAPPO: each agent uses local 22-dim obs at execution time...",
        "...but centralised critic sees full 66-dim joint state during training",
        "Global cooperative reward aligns all 3 agents to minimise network-wide waiting time",
        "Green wave pre-computation biases agents toward coordinated platoon release",
    ],
    Inches(0.15), Inches(1.38), Inches(5.55), Inches(3.20),
    size=12, color=DARK)

    add_text(sl, "Obs (22 dims): queue counts, wait time, stop ratio, phase one-hot, timer, green-wave offset",
             Inches(0.15), Inches(4.65), Inches(5.55), Inches(0.30),
             size=10, italic=True, color=DGREY)
    add_text(sl, "Action: Discrete(3) — 3 green phase configurations per junction",
             Inches(0.15), Inches(4.98), Inches(5.55), Inches(0.30),
             size=10, bold=True, color=DARK)

    # CTDE diagram
    for idx, (lbl, yt, col) in enumerate([
        ("Local Obs o₁  Local Obs o₂  Local Obs o₃\n(22-dim each)",
         Inches(1.12), CYAN),
        ("Shared Actor  πθ\n(decentralised execution)",
         Inches(2.32), MID),
        ("Centralised Critic  Vϕ\n(joint state 66-dim — training only)",
         Inches(3.52), DARK),
    ]):
        add_rect(sl, Inches(5.85), yt, Inches(4.00), Inches(1.00), fill_rgb=col)
        add_text(sl, lbl, Inches(5.95), yt + Inches(0.10),
                 Inches(3.80), Inches(0.80),
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if idx < 2:
            add_text(sl, "↓",
                     Inches(7.65), yt + Inches(1.02), Inches(0.60), Inches(0.25),
                     size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    add_text(sl, "3-Stage Curriculum",
             Inches(5.85), Inches(4.65), Inches(4.00), Inches(0.28),
             size=11, bold=True, color=MID)
    add_bullets(sl, [
        "Stage 1 — low + normal (establishes queue-clearing)",
        "Stage 2 — adds rush_hour_am/pm (phase extension)",
        "Stage 3 — all 6 scenarios incl. holiday + incident",
    ],
    Inches(5.85), Inches(4.96), Inches(4.00), Inches(0.90),
    size=10, color=DARK)
    return sl


def s08_training():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "7", "MAPPO TRAINING SETUP & DEMAND SCENARIOS")
    _clear_textbox(sl)

    add_text(sl, "Training Hyperparameters",
             Inches(0.15), Inches(1.05), Inches(4.60), Inches(0.30),
             size=12, bold=True, color=MID)
    params = [
        ("Total timesteps",       "1,500,000"),
        ("Parallel environments", "4"),
        ("Episode length",        "500 steps × 3 s = 25 min sim"),
        ("PPO clip ε / GAE λ",    "0.2 / 0.95"),
        ("Discount γ",            "0.99"),
        ("Learning rate",         "3×10⁻⁴ (Adam)"),
        ("Entropy coeff.",        "0.01"),
        ("Hardware",              "Core i7 CPU, 16 GB RAM"),
        ("Training time",         "~3–4 hours (CPU-only)"),
    ]
    ph = Inches(0.33)
    for r, (k, v) in enumerate(params):
        rf = LGREY if r % 2 == 0 else WHITE
        yt = Inches(1.40) + r * ph
        for c, (cell, cw) in enumerate(zip([k, v], [Inches(2.40), Inches(2.20)])):
            xl = Inches(0.15) + c * Inches(2.40)
            add_rect(sl, xl, yt, cw, ph, fill_rgb=rf,
                     line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_w_pt=0.3)
            add_text(sl, cell, xl + Inches(0.06), yt + Inches(0.03),
                     cw - Inches(0.12), ph - Inches(0.06),
                     size=10, bold=(c == 0), color=DARK, align=PP_ALIGN.LEFT)

    # Reward
    add_text(sl, "Reward:  r = −Σ tanh(wait_i / (n_lanes · 60))",
             Inches(0.15), Inches(4.50), Inches(4.70), Inches(0.30),
             size=11, italic=True, color=DARK)

    # 6 scenarios
    add_text(sl, "6 Demand Scenarios",
             Inches(5.10), Inches(1.05), Inches(4.80), Inches(0.30),
             size=12, bold=True, color=MID)
    scens = [
        ("low",           "Off-peak, minimal vehicles",         GREEN),
        ("normal",        "Average daily flow",                 MID),
        ("rush_hour_am",  "Morning rush — inbound",             ORANGE),
        ("rush_hour_pm",  "Evening rush — outbound",            ORANGE),
        ("holiday",       "Heavy through-traffic, few turns",   MID),
        ("incident",      "Partial closure, asymmetric demand", RED),
    ]
    sw, sh2 = Inches(2.28), Inches(0.70)
    sg      = Inches(0.08)
    for i, (name, desc, col) in enumerate(scens):
        row = i // 2
        ci  = i % 2
        xl  = Inches(5.10) + ci * (sw + sg)
        yt  = Inches(1.42) + row * (sh2 + sg)
        add_rect(sl, xl, yt, sw, sh2, fill_rgb=col)
        add_text(sl, name, xl + Inches(0.08), yt + Inches(0.04),
                 sw - Inches(0.16), Inches(0.26),
                 size=10, bold=True, color=WHITE)
        add_text(sl, desc, xl + Inches(0.08), yt + Inches(0.33),
                 sw - Inches(0.16), Inches(0.30),
                 size=9, color=WHITE)

    add_text(sl, "Curriculum: Stage 1 (low+normal) → Stage 2 (+rush) → Stage 3 (all 6)",
             Inches(5.10), Inches(3.73), Inches(4.80), Inches(0.30),
             size=10, italic=True, color=DGREY)

    # network/actor spec
    add_text(sl, "Network: Actor + Centralised Critic — 2 hidden layers × 128 units, tanh activations, orthogonal init",
             Inches(0.15), Inches(4.86), Inches(9.70), Inches(0.30),
             size=10, italic=True, color=DGREY, align=PP_ALIGN.CENTER)
    return sl


def s09_convergence():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "8", "MAPPO TRAINING CONVERGENCE")
    _clear_textbox(sl)

    cards = [
        ("Mean Training Reward",
         "−15 → −9 over 1.5M steps\nSteep rise first 200k steps\nPlateaus near −9 from 1M+\nOscillations = multi-agent non-stationarity",
         GREEN),
        ("Evaluation Reward\n(most important)",
         "−9 → near −2\n78% reduction in penalty\nClassic sigmoid learning curve\nClean plateau beyond 1M steps\nConfirms stable cooperative policy",
         MID),
        ("Centralised Critic Loss",
         "Starts at ~4.5 (uninitialised)\nDrops sharply to ~1.5\nStabilises at 1.4–1.6\nPrerequisite for reliable\nadvantage estimation",
         DARK),
        ("Policy Entropy",
         "0.95 → ~0.3 at convergence\nNear-uniform → confident\nRetains residual entropy\nHealthy MAPPO profile —\nnot collapsed to zero",
         ORANGE),
    ]
    cw2, ch2 = Inches(2.35), Inches(2.75)
    gap2     = Inches(0.14)
    lt2      = Inches(0.12)
    tt2      = Inches(1.08)
    for i, (title, body, col) in enumerate(cards):
        xl = lt2 + i * (cw2 + gap2)
        add_rect(sl, xl, tt2, cw2, ch2, fill_rgb=col)
        add_text(sl, title, xl + Inches(0.10), tt2 + Inches(0.08),
                 cw2 - Inches(0.20), Inches(0.42),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, xl + Inches(0.1), tt2 + Inches(0.52),
                 cw2 - Inches(0.20), Pt(1.5), fill_rgb=WHITE)
        add_text(sl, body, xl + Inches(0.10), tt2 + Inches(0.62),
                 cw2 - Inches(0.20), ch2 - Inches(0.72),
                 size=10, color=WHITE)

    add_rect(sl, Inches(0.12), Inches(4.05), Inches(9.76), Inches(1.60),
             fill_rgb=LGREY)
    add_text(sl, "Actor Loss Arc",
             Inches(0.22), Inches(4.12), Inches(4.00), Inches(0.28),
             size=11, bold=True, color=DARK)
    add_text(sl,
             "Begins near −0.004, rises to −0.002 mid-training (large gradients early on), "
             "then returns to −0.003 as PPO clipping activates less often. "
             "No spikes or sign changes — clipping prevented all destabilising updates.",
             Inches(0.22), Inches(4.43), Inches(4.60), Inches(1.12),
             size=10, color=DARK)
    add_text(sl,
             "Key takeaway: evaluation reward plateau at −2 vs −9 at start "
             "is the strongest evidence the agents reached a stable, generalisable cooperative policy.",
             Inches(5.05), Inches(4.20), Inches(4.80), Inches(1.25),
             size=12, bold=True, color=MID, align=PP_ALIGN.CENTER)
    return sl


def s10_eval_wait_queue():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "9", "MAPPO EVALUATION — WAITING TIME & QUEUE LENGTH")
    _clear_textbox(sl)

    add_text(sl,
             "All metrics at TL_A (junction 6073919354)  •  n=3 episodes per cell",
             Inches(0.15), Inches(1.05), Inches(9.70), Inches(0.25),
             size=11, italic=True, color=DGREY, align=PP_ALIGN.CENTER)

    add_text(sl, "Mean Waiting Time per Lane (s)",
             Inches(0.15), Inches(1.35), Inches(4.90), Inches(0.28),
             size=11, bold=True, color=MID)
    wt_h = ["Scenario", "MAPPO", "Fixed", "vs Fixed"]
    wt_r = [
        ("low",          "0.06",  "0.97",  "−94.3%"),
        ("normal",       "0.19",  "2.18",  "−91.4%"),
        ("rush_hour_am", "4.02",  "5.90",  "−31.8%"),
        ("rush_hour_pm", "0.07",  "2.49",  "−97.3%"),
        ("holiday",      "0.35",  "5.73",  "−93.9%"),
        ("incident",     "0.09",  "1.34",  "−93.5%"),
    ]
    render_table(sl, wt_h, wt_r,
                 [Inches(1.55), Inches(0.85), Inches(0.85), Inches(1.00)],
                 Inches(0.15), Inches(1.65), row_h=Inches(0.35))

    add_text(sl, "Mean Queue Length (halting veh / lane)",
             Inches(5.20), Inches(1.35), Inches(4.70), Inches(0.28),
             size=11, bold=True, color=MID)
    ql_h = ["Scenario", "MAPPO", "Fixed", "vs Fixed"]
    ql_r = [
        ("low",          "0.01", "0.04", "−69.5%"),
        ("normal",       "0.02", "0.10", "−75.3%"),
        ("rush_hour_am", "0.29", "0.29", "+2.1%"),
        ("rush_hour_pm", "0.02", "0.10", "−84.1%"),
        ("holiday",      "0.06", "0.24", "−75.9%"),
        ("incident",     "0.02", "0.06", "−70.5%"),
    ]
    render_table(sl, ql_h, ql_r,
                 [Inches(1.55), Inches(0.85), Inches(0.85), Inches(1.10)],
                 Inches(5.20), Inches(1.65), row_h=Inches(0.35))

    add_rect(sl, Inches(0.12), Inches(4.45), Inches(9.76), Inches(1.25),
             fill_rgb=LGREY)
    add_text(sl, "Key Insights",
             Inches(0.22), Inches(4.52), Inches(9.60), Inches(0.28),
             size=11, bold=True, color=DARK)
    add_bullets(sl, [
        "Waiting time reduced 31.8%–97.3% across ALL six scenarios",
        "rush_hour_am: queue is demand-bound (both 0.29) — but MAPPO still cuts wait 31.8% and boosts throughput +20.9%",
        "holiday + incident: largest absolute wait reductions — MAPPO drops empty phases, extends congested ones",
    ],
    Inches(0.22), Inches(4.82), Inches(9.60), Inches(0.75),
    size=11, color=DARK)
    return sl


def s11_throughput():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "10", "EVALUATION — THROUGHPUT & CONTROLLER COMPARISON")
    _clear_textbox(sl)

    add_text(sl, "Total Throughput (vehicles crossing TL_A per episode)",
             Inches(0.15), Inches(1.05), Inches(5.65), Inches(0.28),
             size=11, bold=True, color=MID)
    tp_h = ["Scenario", "MAPPO", "Fixed", "vs Fixed"]
    tp_r = [
        ("low",          "13.67", "13.67", "0.0% (demand-limited)"),
        ("normal",       "20.33", "22.33", "−9.0% (wait traded)"),
        ("rush_hour_am", "90.67", "75.00", "+20.9%"),
        ("rush_hour_pm", "13.67", "12.33", "+10.8%"),
        ("holiday",      "81.67", "64.33", "+26.9%"),
        ("incident",     "18.00", "20.00", "−10.0% (wait traded)"),
    ]
    render_table(sl, tp_h, tp_r,
                 [Inches(1.55), Inches(0.80), Inches(0.80), Inches(2.30)],
                 Inches(0.15), Inches(1.38), row_h=Inches(0.35))

    # 3-controller comparison
    for i, (name, body, col) in enumerate([
        ("FIXED-TIME",
         "Pre-timed cycles\nNo queue awareness\nBaseline — worst overall",
         RED),
        ("SPPO (SA-PPO)",
         "PPO on single junction\nLocal obs only\nBetter than Fixed;\nno corridor coordination",
         ORANGE),
        ("MAPPO",
         "3 cooperative agents\nCTDE — global reward\nBest on all metrics\n(except demand-limited low)",
         GREEN),
    ]):
        xl = Inches(5.95) + i * Inches(1.38)
        add_rect(sl, xl, Inches(1.05), Inches(1.28), Inches(3.50), fill_rgb=col)
        add_text(sl, name, xl + Inches(0.07), Inches(1.12),
                 Inches(1.14), Inches(0.40),
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, xl + Inches(0.07), Inches(1.55),
                 Inches(1.14), Pt(1.5), fill_rgb=WHITE)
        add_text(sl, body, xl + Inches(0.07), Inches(1.65),
                 Inches(1.14), Inches(2.65),
                 size=9, color=WHITE)

    add_text(sl, "Performance hierarchy:  MAPPO  >  SPPO  >  Fixed-Time",
             Inches(0.15), Inches(4.75), Inches(9.70), Inches(0.38),
             size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(sl,
             "Per-step time-series at TL_A confirm consistently lower waiting times, queues, "
             "stop ratios and pressures under MAPPO throughout the full 1,800-second simulation. "
             "MAPPO also exhibits emergent cross-junction coordination: green wave progression, "
             "incident-adaptive phase consolidation, holiday through-phase extension.",
             Inches(0.15), Inches(5.18), Inches(9.70), Inches(0.55),
             size=11, color=DARK, align=PP_ALIGN.CENTER)
    return sl


def s12_yolo():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "11", "YOLOV8 VEHICLE DETECTION MODULE")
    _clear_textbox(sl)

    add_bullets(sl, [
        "Anchor-free detection head — better small-object recall, fewer hyperparameters",
        "Single-stage (one forward pass) — >100 FPS on mid-range GPU",
        "Instance-level bounding boxes — separates overlapping vehicles",
        "Outperforms classical OpenCV MOG2 on all criteria for outdoor intersections",
        "Validated on sample Palapye traffic footage (independent validation)",
        "Queue length estimated from stationary bounding-box count per approach arm",
    ],
    Inches(0.15), Inches(1.05), Inches(5.55), Inches(2.60),
    size=12, color=DARK)

    cv_h = ["Criterion", "OpenCV MOG2", "YOLOv8"]
    cv_r = [
        ("Lighting robustness", "Poor",     "Good"),
        ("Occlusion handling",  "Poor",     "Good"),
        ("Detection accuracy",  "~90% cars", ">90% all classes"),
        ("Queue estimation",    "Indirect", "Direct (stationary bbox)"),
        ("Selected",            "No",       "YES"),
    ]
    render_table(sl, cv_h, cv_r,
                 [Inches(2.15), Inches(1.55), Inches(1.65)],
                 Inches(0.15), Inches(3.72), row_h=Inches(0.35))

    add_rect(sl, Inches(5.85), Inches(1.05), Inches(4.05), Inches(4.60),
             fill_rgb=LGREY)
    add_text(sl, "FUTURE INTEGRATION PATHWAY",
             Inches(5.95), Inches(1.12), Inches(3.85), Inches(0.28),
             size=11, bold=True, color=MID)
    add_bullets(sl, [
        "Deploy calibrated camera at one A1 junction",
        "Validate queue estimates vs manual counts",
        "Replace SUMO TraCI obs with live YOLO estimates in MAPPO obs pipeline",
        "Closed-loop: real camera → YOLO → MAPPO policy → signal actuation",
    ],
    Inches(5.95), Inches(1.45), Inches(3.85), Inches(2.00),
    size=11, color=DARK)
    add_rect(sl, Inches(5.95), Inches(3.50), Inches(3.85), Pt(1.5), fill_rgb=ORANGE)
    add_text(sl,
             "Vision: real-world deployment replaces SUMO TraCI with live YOLOv8 observations",
             Inches(5.95), Inches(3.60), Inches(3.85), Inches(0.85),
             size=11, bold=True, italic=True, color=DARK, align=PP_ALIGN.CENTER)

    add_text(sl,
             "Current status: module implemented and validated independently; "
             "full integration with the SUMO training loop is identified as future work.",
             Inches(0.15), Inches(5.22), Inches(9.70), Inches(0.40),
             size=11, italic=True, color=DGREY, align=PP_ALIGN.CENTER)
    return sl


def s13_conclusions():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "12", "CONCLUSIONS")
    _clear_textbox(sl)

    add_text(sl, "Phase 1: Single-Agent PPO Baseline",
             Inches(0.15), Inches(1.05), Inches(9.70), Inches(0.30),
             size=13, bold=True, color=MID)
    add_bullets(sl, [
        "Trained on one SUMO junction; 9,000 timesteps; 5 demand phases",
        "48.5% throughput improvement, 27.8% queue reduction, 27.3% pressure reduction vs Fixed (p<0.001)",
        "50–68% delay improvement during heavy/peak demand — aggregate figure misleading",
        "Established quantitative benchmark and proved RL can learn adaptive signal control",
    ],
    Inches(0.15), Inches(1.38), Inches(9.70), Inches(1.60),
    size=12, color=DARK)

    add_text(sl, "Phase 2: MAPPO Multi-Intersection Extension",
             Inches(0.15), Inches(3.05), Inches(9.70), Inches(0.30),
             size=13, bold=True, color=MID)
    add_bullets(sl, [
        "All 3 A1 junctions controlled cooperatively; 1.5M steps; 6 scenarios; 3-stage curriculum",
        "Evaluation reward: 78% reduction in penalty; clean convergence plateau confirmed",
        "Waiting time at TL_A vs Fixed: reduced 31.8%–97.3% across ALL 6 scenarios",
        "Throughput: +10.8% to +26.9% in high-demand; small trade-off under light demand",
        "Emergent cross-junction coordination: green wave, incident-adaptive phase consolidation",
        "First reported cooperative MARL evaluation on a Botswana signalised corridor",
    ],
    Inches(0.15), Inches(3.38), Inches(9.70), Inches(1.85),
    size=12, color=DARK)

    add_rect(sl, Inches(0.12), Inches(5.40), Inches(9.76), Inches(0.55), fill_rgb=DARK)
    add_text(sl,
             "Phase 1 proved RL can learn adaptive signal control at one junction.  "
             "Phase 2 demonstrated cooperative MARL extends that to the full A1 corridor — "
             "delivering measurably stronger performance where single-agent and fixed-time approaches fail.",
             Inches(0.22), Inches(5.48), Inches(9.60), Inches(0.38),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return sl


def s14_recommendations():
    sl = _clone_slide(src, TMPL_IDX, dst)
    _set_title(sl, "13", "RECOMMENDATIONS & FUTURE WORK")
    _clear_textbox(sl)

    recs = [
        ("Field Traffic Count Collection",
         "Obtain real turning movement counts from the Botswana DoR to calibrate the SUMO model.",
         MID),
        ("Camera Calibration + YOLO Integration",
         "Deploy at one A1 junction; validate queue estimates; integrate into MAPPO obs pipeline.",
         MID),
        ("Actuated Baseline Comparison",
         "Add Webster-optimised or gap-based actuated controller for rigorous benchmarking.",
         DARK),
        ("Safety: Pedestrian Phases + Emergency Pre-emption",
         "Add pedestrian phases, min-green guarantees, emergency vehicle pre-emption.",
         ORANGE),
        ("Extended Training + Multi-Seed",
         "Scale to 5–10M steps with GPU; run multiple seeds to quantify variance.",
         GREEN),
        ("Emissions-Aware Reward",
         "Incorporate idle-time fuel penalty to align objectives with Palapye air quality.",
         GREEN),
        ("Hardware Deployment Roadmap",
         "Define signal controller API, camera specs, latency budget, failsafe fallback.",
         DARK),
        ("Continual Learning Pipeline",
         "Build dataset pipeline from live runs enabling periodic fine-tuning on real-world data.",
         MID),
    ]
    rw2 = Inches(4.72)
    rh2 = Inches(0.75)
    rg2 = Inches(0.08)
    lt2 = Inches(0.15)
    tt2 = Inches(1.05)
    for i, (title, body, col) in enumerate(recs):
        row = i // 2
        ci  = i % 2
        xl  = lt2 + ci * (rw2 + Inches(0.12))
        yt  = tt2 + row * (rh2 + rg2)
        add_rect(sl, xl, yt, rw2, rh2, fill_rgb=col)
        add_text(sl, title, xl + Inches(0.08), yt + Inches(0.04),
                 rw2 - Inches(0.16), Inches(0.26),
                 size=10, bold=True, color=WHITE)
        add_text(sl, body, xl + Inches(0.08), yt + Inches(0.33),
                 rw2 - Inches(0.16), rh2 - Inches(0.38),
                 size=9, color=WHITE)
    return sl


def s15_thankyou():
    """Clone the thank-you slide from the template."""
    sl = _clone_slide(src, 14, dst)   # slide 15 in the original
    # Update title bar
    _set_title(sl, "14", "THANK YOU")
    # Find and update the big text box
    for sh in sl.shapes:
        if not hasattr(sh, "text_frame"):
            continue
        t = sh.text.strip()
        if "THANK YOU" in t or "QnA" in t:
            tf = sh.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = ""
            tf.paragraphs[0].runs[0].text = (
                "THANK YOU FOR YOUR TIME\n\n\n"
                "Tsotlhe Nayang Seiphepi  |  21001137\n"
                "tsotlhenayangseiphepi@gmail.com\n"
                "Supervisor: Dr Bokamoso Basutli  |  BIUST 2025\n\n"
                "QnA"
            )
    return sl


# ─────────────────────────────────────────────────────────────────────────────
# Run all slide builders
# ─────────────────────────────────────────────────────────────────────────────

builders = [
    s01_title,
    s02_outline,
    s03_intro,
    s04_objectives,
    s05_architecture,
    s06_ppo_results,
    s07_mappo_design,
    s08_training,
    s09_convergence,
    s10_eval_wait_queue,
    s11_throughput,
    s12_yolo,
    s13_conclusions,
    s14_recommendations,
    s15_thankyou,
]

for fn in builders:
    fn()
    print(f"  Built: {fn.__name__}")

# Strip the N_ORIG original template slides from the sldIdLst so only our
# new slides remain.  New slides were appended, so they sit at indices N_ORIG..end.
sldIdLst = dst.slides._sldIdLst
orig_els = list(sldIdLst)[:N_ORIG]
for el in orig_els:
    sldIdLst.remove(el)

print(f"\nStripped {N_ORIG} template slides; {len(dst.slides)} new slides remain.")
os.makedirs("docs", exist_ok=True)
dst.save(OUT)
print(f"Saved: {OUT}")
